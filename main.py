import os#подключаем библиотеку os
from pptx import Presentation#ипортируем и подключаем из библиотеки pptx класс Presentation
FILE_NAME='Osennyaya_igra_3.pptx'#записываем в переменную FILE_NAME путь к файлу 'Osennyaya_igra_3.pptx'

def is_not_valid(text: str) -> bool:
    return ' ' in text or '-' in text or ':' in text or 'СУПЕРКРОКО' in text

words = list()

current_file=os.path.realpath(__file__)#в переменную записываем путь к текущему файлу
print(f'{current_file=}')#выводим путь текущего файла
current_directory=os.path.dirname(current_file)#в переменную записываем путь к текущей директории
print(f'{current_directory=}')#выводим путь текущей директории

prs = Presentation(FILE_NAME)#записываем в переменную данные из файла
for slide in prs.slides:#перебираем слайды
    for shape in slide.shapes:#перебираем форму
        if not shape.has_text_frame:#условие, чтобы на слайде не было текста
            continue#если условие выполняется, продолжаем перебирать слайды
        if is_not_valid(shape.text):    
            continue
        words.append(shape.text)#если условие не выполняется, то выводим текст со слайда 


with open("words.txt", "w", encoding="utf-8") as file:
    print(*words, file=file)