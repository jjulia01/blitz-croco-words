from pathlib import Path  #импортируем Path из библиотеки pathlib
from zipfile import ZipFile #импортируем ZipFile из библиотеки zipfile 

from pptx import Presentation #импортируем Presentation из библиотеки pptx 
# Константа имени файла
FILENAME_ZIP: str = "croco-blitz-source.zip" 

def read_zip_file(): #создаём функцию read_zip_file для чтения файла
    print(FILENAME_ZIP) #выводим значение переменной FILENAME_ZIP


def present()-> None: #создаем функцию present
    prs:Presentation = Presentation("Osennyaya_igra_3.pptx") #создаем переменную типа Presentation
    for slide in prs.slides: #запускаем цикл for по slide
        for shape in slide.shapes: #запускаем цикл for по shape
            if not shape.has_text_frame: #запускаем условную конструкцию if
                continue #если условие выполняется, то дейстия повторяются начиная со строки for shape in slide.shapes:
            print(shape.text_frame) #если условие не выполняется, то выводим значение shape.text_frame
if __name__ == "__main__": 
    read_zip_file() #выводим значение функции
    present() #выводим значение функции